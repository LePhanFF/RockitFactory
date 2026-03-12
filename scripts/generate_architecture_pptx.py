#!/usr/bin/env python3
"""Generate architecture evolution PowerPoint from brainstorm docs."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_GREEN = RGBColor(0x00, 0xC9, 0x7B)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
SUBTLE_BG = RGBColor(0xF5, 0xF5, 0xF5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=WHITE, bullet_color=ACCENT_BLUE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(4)
        p.level = 0
    return txBox


def add_rounded_box(slide, left, top, width, height, text, fill_color, text_color=WHITE,
                    font_size=11, bold=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = "Calibri"
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_arrow(slide, left, top, width, height, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========== SLIDE 1: Title ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 1, 1.5, 11, 1.5, "RockitFactory", 44, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, 1, 3.0, 11, 1, "Agentic Trading Intelligence Platform", 28, ACCENT_BLUE, False, PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.2, 11, 0.8, "Architecture Evolution: From Deterministic Signals to LLM-Powered Agent Decisions",
                 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    add_text_box(slide, 1, 5.8, 11, 0.5, "NQ Futures  |  272 Sessions  |  Dalton Market Profile  |  Qwen3.5 + DuckDB",
                 13, MID_GRAY, False, PP_ALIGN.CENTER)
    add_text_box(slide, 1, 6.5, 11, 0.5, "March 2026", 14, MID_GRAY, False, PP_ALIGN.CENTER)

    # ========== SLIDE 2: Problem Statement ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.4, 11, 0.8, "The Problem", 32, ACCENT_BLUE, True)
    add_bullet_list(slide, 0.8, 1.3, 5.5, 4.5, [
        "6 separate repositories, no unified pipeline",
        "Strategy signals fire without context — no \"should I take this?\"",
        "No way to correlate trades with market structure at entry",
        "LLM analysis exists but disconnected from decision-making",
        "Manual trade review — no systematic self-improvement loop",
        "Backtesting doesn't include agent intelligence",
    ], 14, LIGHT_GRAY)

    add_text_box(slide, 7, 0.4, 5.5, 0.8, "The Vision", 32, ACCENT_GREEN, True)
    add_bullet_list(slide, 7, 1.3, 5.5, 4.5, [
        "Monorepo: research, backtest, train, serve, trade",
        "Agents evaluate every signal before execution",
        "DuckDB warehouse: trades + deterministic context + LLM analysis",
        "Trained LLM as analyst — your trader voice interpreting data",
        "Self-learning loop: review, observe, adapt, improve",
        "One process, two loops: strategy runner (1-min) + orchestrator (5-min)",
    ], 14, LIGHT_GRAY)

    # ========== SLIDE 3: Three-Tier Model Strategy ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.4, 11, 0.8, "Three-Tier Model Strategy", 32, WHITE, True)

    # Tier 0
    add_rounded_box(slide, 0.8, 1.5, 3.5, 1.8,
                    "TIER 0\nDeterministic Python\n\n38 modules, <10ms\n80% of the work", ACCENT_GREEN, WHITE, 13)
    # Tier 1
    add_rounded_box(slide, 4.8, 1.5, 3.5, 1.8,
                    "TIER 1\nQwen3.5 Local (Ollama/vLLM)\n\nAgent debates, real-time\n~2 min per signal", ACCENT_BLUE, WHITE, 13)
    # Tier 2
    add_rounded_box(slide, 8.8, 1.5, 3.5, 1.8,
                    "TIER 2\nOpus 4.6 / Gemini API\n\nMeta-review, design\nPeriodic, not real-time", ACCENT_ORANGE, WHITE, 13)

    add_bullet_list(slide, 0.8, 3.8, 11.5, 3, [
        "Tier 0: Deterministic modules (IB, TPO, DPOC, CRI, delta, FVG, etc.) — free, instant, always running",
        "Tier 1: Local LLM on DGX Spark (128GB) — Advocate/Skeptic debate on signal, trained analyst on 5-min timer",
        "Tier 2: Frontier API models — meta-reviews every 1-3 days, architecture decisions, deep backtesting analysis",
        "Key principle: LLM = analyst/tape reader, NOT trader. Strategies emit signals, agents evaluate.",
        "Single Qwen3.5 + single LoRA. Agent roles via system prompts, not separate models.",
    ], 13, LIGHT_GRAY)

    # ========== SLIDE 4: Current Architecture ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.3, 11, 0.8, "Current Architecture — Signal to Decision Pipeline", 28, WHITE, True)

    # Pipeline boxes
    y = 1.3
    add_rounded_box(slide, 0.5, y, 2.2, 0.9, "Strategy Signal\n(5 strategies)", RGBColor(0x44, 0x44, 0x66), WHITE, 11)
    add_rounded_box(slide, 3.0, y, 2.2, 0.9, "Mechanical Filters\nBias + DayType + AntiChase", RGBColor(0x44, 0x66, 0x44), WHITE, 10)
    add_rounded_box(slide, 5.5, y, 2.2, 0.9, "CRI Gate +\nObservers (9 cards)", ACCENT_BLUE, WHITE, 11)
    add_rounded_box(slide, 8.0, y, 2.2, 0.9, "DuckDB\nHistorical Stats", RGBColor(0x66, 0x44, 0x88), WHITE, 11)
    add_rounded_box(slide, 10.5, y, 2.2, 0.9, "ADVOCATE\n(Qwen3.5 LLM)", ACCENT_GREEN, WHITE, 11)

    y2 = 2.6
    add_rounded_box(slide, 10.5, y2, 2.2, 0.9, "SKEPTIC\n(Qwen3.5 LLM)", ACCENT_ORANGE, WHITE, 11)
    add_rounded_box(slide, 7.5, y2, 2.5, 0.9, "ORCHESTRATOR\n(Deterministic Scorecard)", ACCENT_RED, WHITE, 11)
    add_rounded_box(slide, 4.5, y2, 2.5, 0.9, "TAKE / SKIP /\nREDUCE_SIZE", RGBColor(0x00, 0x80, 0x60), WHITE, 12, True)

    # Results table
    add_text_box(slide, 0.5, 3.9, 12, 0.6, "Backtest Results — 272 Sessions, NQ Futures", 20, ACCENT_BLUE, True)

    table_data = [
        ["Run", "Configuration", "Trades", "Win Rate", "PF", "Net PnL", "$/Trade"],
        ["A", "No filters (baseline)", "408", "56.1%", "2.45", "$159,332", "$390"],
        ["B", "Mechanical filters only", "259", "61.0%", "3.07", "$125,885", "$486"],
        ["C", "Mechanical + Det. Agent", "205", "64.4%", "3.33", "$99,000", "$483"],
        ["E", "Mech + Agent + LLM Debate", "179", "66.5%", "3.58", "$92,909", "$519"],
    ]
    rows, cols = len(table_data), len(table_data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(4.5), Inches(12), Inches(2.3))
    tbl = tbl_shape.table

    col_widths = [0.6, 3.5, 1.2, 1.2, 1.0, 1.8, 1.2]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w)

    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = table_data[r][c]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER if c > 1 else PP_ALIGN.LEFT
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                elif r == 4:  # Run E highlight
                    p.font.bold = True
                    p.font.color.rgb = ACCENT_GREEN
                else:
                    p.font.color.rgb = DARK_TEXT

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x55)
            elif r == 4:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x2A)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = SUBTLE_BG if r % 2 == 0 else WHITE

    # ========== SLIDE 5: The Courtroom Analogy ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.4, 11, 0.8, "The Courtroom — Agent Pipeline as Legal Process", 28, WHITE, True)

    roles = [
        ("DOMAIN EXPERTS\n(38 Deterministic Modules)", "Specialized factual testimony:\nIB extended 1.8x, 4 wick rejections\nat VAH, DPOC migrating up", ACCENT_BLUE, 0.5),
        ("ANALYST\n(Trained Qwen3.5 — Your Voice)", "Interprets expert data:\n\"P-shape forming, sellers trapped\nabove POC, this is a fade setup\"", ACCENT_GREEN, 3.4),
        ("ADVOCATE LAWYER\n(LLM Agent)", "Argues FOR the trade\nusing expert data +\nanalyst interpretation", ACCENT_ORANGE, 6.3),
        ("DEFENSE LAWYER\n(LLM Agent)", "Challenges the case:\nflags warnings, disputes\nevidence, finds weaknesses", ACCENT_RED, 9.2),
    ]

    for title, desc, color, x in roles:
        add_rounded_box(slide, x, 1.4, 2.6, 1.6, title, color, WHITE, 11)
        add_text_box(slide, x, 3.2, 2.6, 1.5, desc, 10, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    # Judge
    add_rounded_box(slide, 4.5, 5.0, 4.3, 1.0,
                    "JUDGE — Orchestrator (Deterministic Scorecard)\nTAKE  /  SKIP  /  REDUCE_SIZE",
                    RGBColor(0x80, 0x00, 0x80), WHITE, 13)

    add_text_box(slide, 0.5, 6.3, 12, 0.8,
                 "Court convenes ON DEMAND only — when a signal fires or user asks. Analyst runs on 5-min timer (amortized). No signal = no trial.",
                 13, MID_GRAY, False, PP_ALIGN.CENTER)

    # ========== SLIDE 6: Strategy Results ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.4, 11, 0.8, "Run E — Strategy Performance (LLM Debate)", 28, WHITE, True)

    strat_data = [
        ["Strategy", "Trades", "Win Rate", "Net PnL", "Debated", "Take Rate"],
        ["Opening Range Rev", "52", "75.0%", "$50,064", "56", "93%"],
        ["OR Acceptance", "74", "66.2%", "$22,329", "91", "81%"],
        ["20P IB Extension", "32", "56.2%", "$16,542", "36", "89%"],
        ["B-Day", "20", "65.0%", "$5,655", "26", "77%"],
        ["80P Rule", "1", "0.0%", "-$1,682", "2", "50%"],
        ["TOTAL", "179", "66.5%", "$92,909", "211", "—"],
    ]
    rows, cols = len(strat_data), len(strat_data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(1.5), Inches(1.5), Inches(10), Inches(3.5))
    tbl = tbl_shape.table

    col_widths = [2.8, 1.2, 1.2, 1.5, 1.2, 1.2]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w)

    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = strat_data[r][c]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                if r == 0 or r == rows - 1:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                elif strat_data[r][0] == "Opening Range Rev":
                    p.font.color.rgb = ACCENT_GREEN
                    p.font.bold = True
                elif strat_data[r][0] == "80P Rule":
                    p.font.color.rgb = ACCENT_RED
                else:
                    p.font.color.rgb = DARK_TEXT

            if r == 0 or r == rows - 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x55)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = SUBTLE_BG if r % 2 == 0 else WHITE

    add_bullet_list(slide, 1.5, 5.3, 10, 2, [
        "OR Rev is the crown jewel: 75% WR, $50K net, 93% take rate — LLM has high conviction",
        "LLM essentially killed 80P Rule (2 debates, 1 trade, 1 loss) — correct behavior given known issues",
        "OR Acceptance most debated (91 signals) — LLM selective, 81% take rate",
        "Progressive improvement: A(56.1%) → B(61.0%) → C(64.4%) → E(66.5%) — each layer adds value",
    ], 12, LIGHT_GRAY)

    # ========== SLIDE 7: Target Architecture ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.3, 11, 0.8, "Target Architecture — Production Pipeline", 28, WHITE, True)

    # Always running section
    add_text_box(slide, 0.5, 1.1, 6, 0.5, "ALWAYS RUNNING (timer, amortized)", 16, ACCENT_GREEN, True)

    add_rounded_box(slide, 0.5, 1.7, 2.5, 1.2,
                    "1-min Strategy Runner\n5 strategies\nSignal generation", RGBColor(0x33, 0x55, 0x33), WHITE, 10)
    add_rounded_box(slide, 3.3, 1.7, 2.5, 1.2,
                    "5-min Orchestrator\n38 deterministic modules\nContext building", ACCENT_BLUE, WHITE, 10)
    add_rounded_box(slide, 6.1, 1.7, 2.5, 1.2,
                    "5-min LLM Analyst\nTrained Qwen3.5\nYour trader voice", ACCENT_GREEN, WHITE, 10)
    add_rounded_box(slide, 8.9, 1.7, 2.5, 1.2,
                    "DuckDB + GCS\nAnalysis cached\nHistorical context", RGBColor(0x66, 0x44, 0x88), WHITE, 10)

    # On-demand section
    add_text_box(slide, 0.5, 3.3, 6, 0.5, "ON-DEMAND (signal trigger or user query)", 16, ACCENT_ORANGE, True)

    add_rounded_box(slide, 0.5, 3.9, 1.8, 1.0, "Signal fires\nor user asks", ACCENT_ORANGE, WHITE, 10)
    add_rounded_box(slide, 2.6, 3.9, 2.0, 1.0, "Scorecard\n+ Analyst reports\n(last 3-6)", ACCENT_BLUE, WHITE, 10)
    add_rounded_box(slide, 4.9, 3.9, 1.6, 1.0, "Advocate\n(~50s)", ACCENT_GREEN, WHITE, 11)
    add_rounded_box(slide, 6.8, 3.9, 1.6, 1.0, "Skeptic\n(~60s)", ACCENT_ORANGE, WHITE, 11)
    add_rounded_box(slide, 8.7, 3.9, 1.6, 1.0, "Judge\n(<10ms)", ACCENT_RED, WHITE, 11)
    add_rounded_box(slide, 10.6, 3.9, 1.6, 1.0, "TAKE\nSKIP\nREDUCE", RGBColor(0x00, 0x80, 0x60), WHITE, 11, True)

    # Self-learning loop
    add_text_box(slide, 0.5, 5.3, 6, 0.5, "SELF-LEARNING LOOP (post-market)", 16, RGBColor(0x80, 0x00, 0x80), True)

    add_rounded_box(slide, 0.5, 5.9, 2.2, 0.9, "Trade Reviewer\n(Qwen3.5)", RGBColor(0x55, 0x33, 0x66), WHITE, 10)
    add_rounded_box(slide, 3.0, 5.9, 2.2, 0.9, "Observations\n→ DuckDB", RGBColor(0x55, 0x33, 0x66), WHITE, 10)
    add_rounded_box(slide, 5.5, 5.9, 2.2, 0.9, "Meta-Review\n(Opus 4.6, 1-3 days)", RGBColor(0x55, 0x33, 0x66), WHITE, 10)
    add_rounded_box(slide, 8.0, 5.9, 2.2, 0.9, "Prompt/Param Update\n+ A/B Test", RGBColor(0x55, 0x33, 0x66), WHITE, 10)
    add_rounded_box(slide, 10.5, 5.9, 2.2, 0.9, "Auto-Rollback\nif regression", RGBColor(0x55, 0x33, 0x66), WHITE, 10)

    # ========== SLIDE 8: Speed Roadmap ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.4, 11, 0.8, "Speed Optimization Roadmap", 28, WHITE, True)

    speed_data = [
        ["Optimization", "Speedup", "Effort", "Status"],
        ["Ollama → vLLM (concurrent batching)", "2-3x", "Medium", "Priority"],
        ["Skip obvious signals (debate ambiguous only)", "2x", "Low", "Planned"],
        ["Smaller model (8B/14B benchmark)", "3-4x", "Low", "Planned"],
        ["Session chunking (parallel backtests)", "2-3x", "Low", "Planned"],
        ["Combined estimate", "10-20x", "—", "6.5h → 20-40 min"],
    ]
    rows, cols = len(speed_data), len(speed_data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(1.5), Inches(1.5), Inches(10), Inches(3))
    tbl = tbl_shape.table

    col_widths = [4.5, 1.5, 1.5, 2.0]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w)

    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = speed_data[r][c]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                if r == 0 or r == rows - 1:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = DARK_TEXT
            if r == 0 or r == rows - 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x55)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = SUBTLE_BG if r % 2 == 0 else WHITE

    add_bullet_list(slide, 1.5, 4.8, 10, 2.5, [
        "Model benchmarking: Test GLM-4.7, Qwen3:8B/14B, Phi-4 — if smaller model matches decisions, use it",
        "vLLM continuous batching: same weights, concurrent requests, 2-3x throughput",
        "Skip obvious signals: if deterministic score > 0.7, auto-TAKE without LLM debate",
        "Per-signal latency in production: ~2 min (fits in 5-min candle with room to spare)",
        "Training decision gate: Run E → model benchmark → decide if LoRA training is needed",
    ], 12, LIGHT_GRAY)

    # ========== SLIDE 9: Tech Stack ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.4, 11, 0.8, "Tech Stack", 28, WHITE, True)

    left_items = [
        "Core: Python 3.11+, uv workspaces, monorepo",
        "Backtest: 272 sessions, 5 strategies, MAE/MFE, 38 deterministic modules",
        "Database: DuckDB (local-first, structured retrieval, not vector RAG)",
        "LLM Serving: Ollama → vLLM on DGX Spark (128GB)",
        "Model: Qwen3.5:35b-a3b (MoE, 3B active) — 128K context, 8K output",
        "Training: BF16 LoRA r=64 via Unsloth on DGX Spark",
    ]
    right_items = [
        "Agent Framework: Custom pipeline (Gate → Observers → Debate → Orchestrator)",
        "Filters: YAML-driven CompositeFilter chain (bias, day type, anti-chase, agent)",
        "API: FastAPI (rockit-serve), JWT auth, GCS storage",
        "UI: React 19 + TypeScript + Vite + Tailwind (RockitUI dashboard)",
        "Cloud: GCP (Cloud Run, GCS, Vertex AI) — production only",
        "Testing: 732+ tests, baseline comparison, A/B test framework",
    ]

    add_bullet_list(slide, 0.5, 1.3, 5.8, 5.5, left_items, 13, LIGHT_GRAY)
    add_bullet_list(slide, 6.8, 1.3, 5.8, 5.5, right_items, 13, LIGHT_GRAY)

    # ========== SLIDE 10: What's Next ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.4, 11, 0.8, "What's Next", 32, WHITE, True)

    phases = [
        ("NOW", "Run E complete — LLM debate adds +10.4% WR over baseline", ACCENT_GREEN),
        ("NEXT", "Model benchmarking — test smaller/faster models for same quality", ACCENT_BLUE),
        ("THEN", "Train expert witness — LoRA on deterministic → analysis (your voice)", ACCENT_ORANGE),
        ("AFTER", "Inject LLM analysis stream into agent pipeline → Run F backtest", ACCENT_ORANGE),
        ("PROD", "vLLM deployment + real-time pipeline + RockitUI integration", ACCENT_RED),
    ]

    for i, (label, desc, color) in enumerate(phases):
        y = 1.4 + i * 1.1
        add_rounded_box(slide, 0.8, y, 1.5, 0.8, label, color, WHITE, 16, True)
        add_text_box(slide, 2.6, y + 0.15, 9.5, 0.8, desc, 16, LIGHT_GRAY)

    add_text_box(slide, 0.8, 6.5, 11.5, 0.5,
                 "\"Strategies emit signals. Agents evaluate. LLM reads the tape. The court decides.\"",
                 15, MID_GRAY, True, PP_ALIGN.CENTER)

    return prs


if __name__ == "__main__":
    prs = build_presentation()
    out_path = "reports/RockitFactory-Architecture-Evolution.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")
